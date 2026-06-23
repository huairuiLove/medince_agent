"""
Generate MedSafe backend architecture PPT.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # deep navy
CARD_BG = RGBColor(0x25, 0x25, 0x42)       # card navy
ACCENT = RGBColor(0x00, 0xD2, 0xFF)        # cyan
ACCENT2 = RGBColor(0x7B, 0x2F, 0xFF)       # purple
ACCENT3 = RGBColor(0x00, 0xE6, 0x96)       # green
ACCENT4 = RGBColor(0xFF, 0x6B, 0x6B)       # coral
ACCENT5 = RGBColor(0xFF, 0xD9, 0x3D)       # yellow
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xCC)
MED_GRAY = RGBColor(0x66, 0x66, 0x88)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_number(slide, num, total):
    """Add slide number at bottom right."""
    left = prs.slide_width - Inches(0.8)
    top = prs.slide_height - Inches(0.4)
    tf = slide.shapes.add_textbox(left, top, Inches(0.6), Inches(0.3)).text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = MED_GRAY
    p.alignment = PP_ALIGN.RIGHT


def make_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def make_title_style(tf, text, size=28, color=WHITE, bold=True):
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return p


def add_textbox(slide, left, top, width, height, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_arrow(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_chevron(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


TOTAL_SLIDES = 10

# ════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Accent line at top
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.05))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "MedSafe 后端算法架构", size=44, color=WHITE, bold=True)
add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
            "多模型协作 · 影像处理 · 端到端流程全景", size=22, color=ACCENT)

# Subtitle info
add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
            "基于 FastAPI 的多智能体药物安全审查系统  |  MIMIC-III 临床场景", size=16, color=LIGHT_GRAY)

# Bottom decorative bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.05), prs.slide_width, Inches(0.05))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT2; bar.line.fill.background()

add_number(slide, 1, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
# SLIDE 2: 系统整体架构
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
            "系统整体架构总览", size=32, color=WHITE, bold=True)
add_textbox(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
            "三大核心模块 + 两条处理流水线", size=16, color=LIGHT_GRAY)

# ── Three columns ──
col_w = Inches(3.8)
col_gap = Inches(0.3)
start_x = Inches(0.6)

titles = ["文本药物审查流水线", "影像分析流水线", "实时CPOE & 聊天"]
subs = [
    "LLM多智能体协作 + 规则引擎",
    "VLM视觉模型 + 分割模型",
    "处方实时审查 + ReAct对话"
]
colors = [ACCENT, ACCENT3, ACCENT2]

for i in range(3):
    x = start_x + i * (col_w + col_gap)
    box = make_rounded_rect(slide, x, Inches(1.6), col_w, Inches(1.0), fill_color=CARD_BG, border_color=colors[i])
    tf = box.text_frame; tf.word_wrap = True
    make_title_style(tf, titles[i], size=18, color=colors[i])
    p2 = tf.add_paragraph()
    p2.text = subs[i]; p2.font.size = Pt(12); p2.font.color.rgb = LIGHT_GRAY

# Detail boxes
details = [
    ["提取 (Extract)", "规则引擎 (Rule Gate)", "5个专家Agent并行审查", "多轮对抗辩论 (Debate)", "首席审查仲裁", "协调员澄清"],
    ["影像目录发现", "多模型分割 (6个后端)", "VLM影像分析 (Qwen3-VL)", "DeepSeek报告合成", "融合多Agent药物审查"],
    ["药物目录解析", "DDI-BERT相互作用检测", "ReAct状态机 + 工具调用", "Graph RAG知识检索", "SSE流式输出"]
]

for i in range(3):
    x = start_x + i * (col_w + col_gap)
    y = Inches(2.9)
    for j, item in enumerate(details[i]):
        item_h = Inches(0.55)
        box = make_rounded_rect(slide, x, y + j * (item_h + Inches(0.08)), col_w, item_h, fill_color=RGBColor(0x30, 0x30, 0x55), border_color=colors[i])
        tf = box.text_frame; tf.word_wrap = True
        make_title_style(tf, item, size=13, color=WHITE, bold=False)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Bottom note
add_textbox(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
            "所有流程通过 FastAPI REST API 暴露，支持 SSE 流式、JSON持久化、部门权限隔离", size=13, color=MED_GRAY)

add_number(slide, 2, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
# SLIDE 3: 模型全景图 - 四层模型体系
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
            "四层模型体系 & 协作关系", size=32, color=WHITE, bold=True)
add_textbox(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
            "每一层模型各司其职，通过 Orchestrator 统一调度，数据自上而下流转", size=16, color=LIGHT_GRAY)

# ── Layout: badge (left) + content (right), full width ──
badge_w = Inches(2.6)
content_w = Inches(9.0)
card_w = badge_w + content_w  # ~11.6 inches
card_h = Inches(1.1)
start_y = Inches(1.55)
gap = Inches(0.12)

layers = [
    ("1", "主LLM推理层", "DeepSeek-Chat", ACCENT,
     "5个专家Agent并行推理  |  辩论Critic/Moderator  |  首席仲裁 & 协调员澄清  |  OpenAI兼容API → 结构化JSON输出"),
    ("2", "视觉理解层", "Qwen3-VL (百炼)", ACCENT3,
     "最多12张医学影像  |  X-ray / CT / MRI  |  分割标注叠加输入  |  输出: 临床表现, 影像发现, 药物推荐, 诊断, 风险等级"),
    ("3", "多模态融合层", "DeepSeek-Chat (专用实例)", ACCENT2,
     "VLM结果 + 5个Agent审查意见 + 仲裁结论 + 规则引擎输出  →  7节结构化ClinicalReport  |  逐项交叉验证 & 去重合并"),
    ("4", "本地AI基础层", "Med7 / BERT / 分割模型", ACCENT4,
     "Med7 NER: 药物实体抽取  |  DDI-BERT: 相互作用分类  |  Embedding: 语义药物搜索  |  TotalSegmentator / VISTA3D / SAM / BraTS / CXR"),
]

for i, (num, short_name, model_name, color, desc) in enumerate(layers):
    y = start_y + i * (card_h + gap)

    # Badge on the left
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, badge_w, card_h)
    badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.fill.background()
    tf_badge = badge.text_frame; tf_badge.word_wrap = True; tf_badge.paragraphs[0].alignment = PP_ALIGN.CENTER
    make_title_style(tf_badge, f"Layer {num}", size=11, color=DARK_BG)
    p2 = tf_badge.add_paragraph()
    p2.text = short_name; p2.font.size = Pt(14); p2.font.color.rgb = DARK_BG; p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf_badge.add_paragraph()
    p3.text = model_name; p3.font.size = Pt(9); p3.font.color.rgb = RGBColor(0x20, 0x20, 0x40)
    p3.alignment = PP_ALIGN.CENTER

    # Content card
    content = make_rounded_rect(slide, Inches(0.6) + badge_w + Inches(0.1), y, content_w, card_h,
                                fill_color=CARD_BG, border_color=color)
    tf_c = content.text_frame; tf_c.word_wrap = True
    tf_c.paragraphs[0].text = ""
    p = tf_c.paragraphs[0]
    p.text = desc
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    # Vertical down-arrow between layers (except last)
    if i < 3:
        arrow_y = y + card_h
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                       Inches(0.6) + badge_w + Inches(0.1) + content_w // 2 - Inches(0.1),
                                       arrow_y, Inches(0.22), Inches(0.14))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = color; arrow.line.fill.background()

# Bottom note
add_textbox(slide, Inches(0.6), Inches(6.55), Inches(12), Inches(0.4),
            "Orchestrator 通过 ThreadPoolExecutor 并行调用 Layer 1 的 Agent，Layer 4 模型提供辅助特征，Layer 2 独立触发，Layer 3 融合全部输出",
            size=12, color=MED_GRAY)

add_number(slide, 3, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
# SLIDE 4: 多Agent协作机制
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
            "多智能体协作机制 — MDAgents + ClinicalPilot 模式", size=32, color=WHITE, bold=True)
add_textbox(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
            "五大专家并行审查 → 多轮辩论 → 首席仲裁 → 澄清确认", size=16, color=LIGHT_GRAY)

# Agent boxes
agents = [
    ("临床药师", "药物相互作用\n剂量合理性\n治疗方案评估", ACCENT),
    ("内科医师", "病情整体判断\n合并症分析\n治疗优先级", ACCENT),
    ("过敏专科", "过敏禁忌筛查\n交叉过敏分析\n替代药物建议", ACCENT3),
    ("药房库存", "处方药物可用性\n医保替换建议\n成本效益分析", ACCENT4),
    ("专科路由", "科室特殊规则\n老年/孕产用药\n罕见病专科意见", ACCENT2),
]

agent_w = Inches(2.2)
agent_h = Inches(1.4)
gap = Inches(0.2)
total_w = 5 * agent_w + 4 * gap
start_x = (prs.slide_width - total_w) // 2

for i, (name, desc, color) in enumerate(agents):
    x = start_x + i * (agent_w + gap)
    box = make_rounded_rect(slide, x, Inches(1.5), agent_w, agent_h, fill_color=CARD_BG, border_color=color)
    tf = box.text_frame; tf.word_wrap = True
    make_title_style(tf, name, size=16, color=color)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

# Debate flow
add_textbox(slide, Inches(0.6), Inches(3.2), Inches(12), Inches(0.4),
            "辩论流程 (最高3轮)", size=18, color=ACCENT5, bold=True)
steps = [
    ("Round 1", "所有Agent独立审查\nThreadPool并行"),
    ("Critic分析", "识别分歧点/证据缺口\n安全遗漏/低置信度"),
    ("Round 2-N", "Agent收到Critic反馈\n针对性修订意见"),
    ("停止条件", "达成共识 &\n置信度 ≥ 0.75"),
    ("未达成", "标记人工审核\n安全兜底"),
]

step_w = Inches(2.2)
step_h = Inches(1.2)
total_w = 5 * step_w + 4 * gap
start_x = (prs.slide_width - total_w) // 2

for i, (title, desc) in enumerate(steps):
    x = start_x + i * (step_w + gap)
    color = ACCENT if i < 3 else (ACCENT3 if i == 3 else ACCENT4)
    box = make_rounded_rect(slide, x, Inches(3.7), step_w, step_h, fill_color=CARD_BG, border_color=color)
    tf = box.text_frame; tf.word_wrap = True
    make_title_style(tf, title, size=15, color=color)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

# Arrow connectors between steps
for i in range(4):
    x = start_x + (i+1) * (step_w + gap) - gap - Inches(0.05)
    arr = add_arrow(slide, x, Inches(4.2), Inches(0.25), Inches(0.2), color=MED_GRAY)

# Arbitration & Clarify
add_textbox(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
            "仲裁与澄清", size=18, color=ACCENT5, bold=True)
sections = [
    ("Moderator 合成", "汇总所有轮次意见\n生成结构化辩论总结", ACCENT),
    ("Chief Reviewer 仲裁", "规则优先+LLM意见融合\nrule_strict=true时规则不可覆盖\n安全面板建议也可强制阻断", ACCENT3),
    ("Coordinator 澄清", "生成针对性追问\n补充缺失信息\n兜底确定性澄清引擎", ACCENT2),
]
sect_w = Inches(3.8)
for i, (title, desc, color) in enumerate(sections):
    x = start_x + i * (sect_w + gap + Inches(0.6))
    box = make_rounded_rect(slide, x, Inches(5.6), sect_w, Inches(1.3), fill_color=CARD_BG, border_color=color)
    tf = box.text_frame; tf.word_wrap = True
    make_title_style(tf, title, size=15, color=color)
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY

add_number(slide, 4, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
# SLIDE 5: 规则引擎 (安全兜底)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
            "确定性规则引擎 — 不可被LLM覆盖的安全底座", size=32, color=WHITE, bold=True)
add_textbox(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
            "六维审查 + 规则优先机制", size=16, color=LIGHT_GRAY)

checks = [
    ("药物相互作用", ACCENT, [
        "知识库规则匹配 (drug_pair → risk)",
        "DDI-BERT 模型分类 (Bio_ClinicalBERT)",
        "严重程度分级 → 阻断/警示/安全"
    ]),
    ("重复成分检测", ACCENT2, [
        "相同活性成分多药识别",
        "成分-商品名映射表",
        "防止重复给药过量"
    ]),
    ("特殊人群规则", ACCENT3, [
        "孕期/哺乳期禁忌",
        "儿童/老年人剂量调整",
        "肝肾功能不全药物调整",
        "FDA妊娠分级 D/X 级阻断"
    ]),
    ("过敏禁忌筛查", ACCENT4, [
        "药物-过敏原交叉匹配",
        "同类药物交叉过敏推断",
        "过敏严重程度分层"
    ]),
    ("临床场景审查", ACCENT5, [
        "多重用药检测 (>5种→标黄, >10种→标红)",
        "跌倒风险评估 (Beers标准)",
        "抗胆碱能负担计算"
    ]),
    ("科室优先级规则", ACCENT, [
        "科室特定禁忌/警示",
        "紧急手术消炎优先级",
        "慢病长期用药管理"
    ]),
]

card_w = Inches(3.8)
card_h = Inches(2.0)
for idx, (title, color, items) in enumerate(checks):
    col = idx % 3
    row = idx // 3
    x = Inches(0.6) + col * (card_w + Inches(0.3))
    y = Inches(1.5) + row * (card_h + Inches(0.15))

    box = make_rounded_rect(slide, x, y, card_w, card_h, fill_color=CARD_BG, border_color=color)
    tf = box.text_frame; tf.word_wrap = True
    make_title_style(tf, title, size=16, color=color)
    for item in items:
        p = tf.add_paragraph()
        p.text = f"  • {item}"
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY

# Bottom rule-strict mechanism
y_bot = Inches(5.8)
box = make_rounded_rect(slide, Inches(0.6), y_bot, Inches(12), Inches(1.2), fill_color=RGBColor(0x35, 0x20, 0x40), border_color=ACCENT4)
tf = box.text_frame; tf.word_wrap = True
make_title_style(tf, "🔒 规则优先机制 (rule_strict=true)", size=18, color=ACCENT4)
p = tf.add_paragraph()
p.text = '• 规则引擎输出作为\u201c硬证据\u201d，LLM Agent 不可覆盖                        • DDI-BERT 确认高风险 \u2192 强制阻断'
p.font.size = Pt(13); p.font.color.rgb = WHITE
p = tf.add_paragraph()
p.text = '• 安全面板 (SafetyPanel) 独立运行，与辩论并行 \u2192 最终仲裁时优先采纳            • 过敏专科/临床药师阻断 \u2192 全局阻断'
p.font.size = Pt(13); p.font.color.rgb = WHITE

add_number(slide, 5, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
# SLIDE 6: 影像处理流水线 (1) - 全流程概览
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
            "影像处理流水线 (上) — 从数据源到VLM分析", size=32, color=WHITE, bold=True)
add_textbox(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
            "数据源发现 → 分割标注 → VLM视觉分析 → 报告合成", size=16, color=LIGHT_GRAY)

# Pipeline steps
steps_img = [
    ("① 数据源发现", ACCENT,
     ["ImagingCatalog 扫描5大数据源",
      "MIMIC-CXR (胸片) / MIMIC-CT",
      "BraTS2024 (脑肿瘤MRI)",
      "KiTS19 (肾脏CT)",
      "MONAI-COPD (胸部CT)",
      "→ 生成轴向预览PNG"]),
    ("② 分割模型", ACCENT2,
     ["SerialSegmentService 串行调度",
      "6个后端按需加载/释放GPU",
      "• CXR-Lesion (肺炎病灶)",
      "• BraTS-Tumor (脑肿瘤)",
      "• TotalSegmentator (多器官)",
      "• VISTA3D / SAM-Med3D / SAM2D",
      "→ 每次只一个模型驻留显存"]),
    ("③ 权限控制", ACCENT4,
     ["ImagingScope 部门级授权",
      "每个科室配置可访问的影像源",
      "filter_models() / filter_studies()",
      "path_allowed_for_sources()",
      "→ 用户只能看到/分割授权数据"]),
    ("④ VLM分析", ACCENT3,
     ["OpenAIVisionClient → Qwen3-VL",
      "base64编码图片(最多12张)+临床摘要",
      "支持分割标注覆盖图像输入",
      "多模态提示 → JSON结构化输出",
      "临床表现/影像发现/药物推荐/诊断",
      "过敏/症状/主诉/麻醉/风险等级"]),
]

card_w = Inches(5.8)
card_h = Inches(2.6)
for idx, (title, color, items) in enumerate(steps_img):
    col = idx % 2
    row = idx // 2
    x = Inches(0.6) + col * (card_w + Inches(0.3))
    y = Inches(1.4) + row * (card_h + Inches(0.15))

    box = make_rounded_rect(slide, x, y, card_w, card_h, fill_color=CARD_BG, border_color=color)
    tf = box.text_frame; tf.word_wrap = True
    make_title_style(tf, title, size=18, color=color)
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY if not item.startswith("→") else WHITE
        p.font.bold = item.startswith("→")

# Arrows between steps
arrows = [
    (Inches(6.7), Inches(2.5), Inches(0.3), Inches(0.2), Inches(0.25)),
    (Inches(6.7), Inches(5.1), Inches(0.3), Inches(0.2), Inches(0.25)),
]
for left, top, w, h, _ in arrows:
    add_arrow(slide, left, top, w, h, color=MED_GRAY)

add_number(slide, 6, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
# SLIDE 7: 影像处理流水线 (2) - VLM分析详解
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
            "影像处理流水线 (下) — VLM分析 & DeepSeek报告合成", size=32, color=WHITE, bold=True)

# Left: VLM detail
box1 = make_rounded_rect(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(5.8), fill_color=CARD_BG, border_color=ACCENT3)
tf1 = box1.text_frame; tf1.word_wrap = True
make_title_style(tf1, "Qwen3-VL 视觉语言模型分析", size=20, color=ACCENT3)

items_vlm = [
    "",
    "▸ 输入构造：",
    "   1. 图片 base64 编码 → data:image/png;base64,...",
    "   2. 构建多模态 messages: [system + user(content=[text, image_urls...])]",
    "   3. System prompt 包含：患者摘要、影像模态、分析任务",
    "",
    "▸ API 调用：",
    "   POST 百炼 Model Studio (OpenAI兼容)",
    "   model: qwen3-vl-plus",
    "   max_tokens: 16000",
    "",
    "▸ 结构化 JSON 输出字段：",
    "   • clinical_analysis     — 临床综合分析",
    "   • imaging_findings       — 影像学发现",
    "   • medication_recommendation — 用药建议",
    "   • recommended_drugs      — 推荐药物列表",
    "   • allergies / diagnoses / symptoms — 过敏/诊断/症状",
    "   • chief_complaint / anesthesia_surgery — 主诉/麻醉手术",
    "   • reasoning / risk_level — 推理过程/风险等级",
]
for item in items_vlm:
    p = tf1.add_paragraph()
    p.text = item
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE if item.startswith("▸") else (ACCENT if item.startswith("   •") else LIGHT_GRAY)

# Right: DeepSeek synthesis
box2 = make_rounded_rect(slide, Inches(6.7), Inches(1.2), Inches(5.8), Inches(5.8), fill_color=CARD_BG, border_color=ACCENT2)
tf2 = box2.text_frame; tf2.word_wrap = True
make_title_style(tf2, "DeepSeek 融合报告合成", size=20, color=ACCENT2)

items_syn = [
    "",
    "▸ 输入融合：",
    "   1. Qwen3-VL 视觉分析结果 (JSON)",
    "   2. 5个专家Agent审查意见 (JSON)",
    "   3. 首席审查仲裁结论 (JSON)",
    "   4. 规则引擎审查输出 (JSON)",
    "",
    "▸ 合成策略：",
    "   model: deepseek-chat (专用独立实例)",
    "   system prompt 设定为'资深临床审查专家'",
    "   要求逐项对比、交叉验证、去重合并",
    "",
    "▸ 输出 ClinicalReport (7节)：",
    "   ① clinical_analysis      — 综合分析",
    "   ② imaging_findings        — 影像发现",
    "   ③ medication_recommendation — 用药建议",
    "   ④ pharmacy_assessment     — 药学评估",
    "   ⑤ allergy_analysis       — 过敏分析",
    "   ⑥ anesthesia_surgery     — 麻醉手术",
    "   ⑦ risk_summary + chain_of_thought — 风险总结+推理链",
]
for item in items_syn:
    p = tf2.add_paragraph()
    p.text = item
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE if item.startswith("▸") else (ACCENT if item.startswith("   ") else LIGHT_GRAY)

add_number(slide, 7, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
# SLIDE 8: 完整请求生命周期
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
            "完整请求生命周期 — POST /api/v1/multi-consult", size=30, color=WHITE, bold=True)

# Flow steps as horizontal pipeline
flow_steps = [
    ("1\n提取", ACCENT),
    ("2\n规则门", ACCENT4),
    ("3\nAgent选择", ACCENT2),
    ("4\n并行审查", ACCENT),
    ("5\n多轮辩论", ACCENT3),
    ("6\n安全面板", ACCENT4),
    ("7\nModerator", ACCENT2),
    ("8\n首席仲裁", ACCENT),
    ("9\n澄清", ACCENT3),
    ("10\n持久化", ACCENT5),
]

step_w = Inches(1.15)
step_h = Inches(0.7)
total_w = 10 * step_w + 9 * Inches(0.1)
start_x = (prs.slide_width - total_w) // 2 - Inches(0.1)

for i, (name, color) in enumerate(flow_steps):
    x = start_x + i * (step_w + Inches(0.1))
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, Inches(1.5), step_w, step_h)
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    make_title_style(tf, name, size=11, color=DARK_BG)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Detail cards below the flow
detail_cards = [
    ("步骤 1-2: 预处理阶段", ACCENT,
     "ExtractAgent 调用主LLM解析临床文本 → ExtractionOutput (年龄/性别/症状/诊断/用药/过敏/孕期等)\n"
     "ReviewEngine 确定性六维审查 → 规则输出 (风险等级/证据/阻断/澄清目标)\n"
     "Med7 NER 辅助药物实体抽取 → 补充LLM提取结果"),
    ("步骤 3-5: 多Agent协作阶段", ACCENT3,
     "Orchestrator._active_agents() → 必选4个Agent + 条件触发的SpecialistAgent/DepartmentSpecialist\n"
     "ThreadPoolExecutor 并行调用所有Agent → 每个独立调用主LLM获得审查意见\n"
     "CriticAgent 分析分歧/证据缺口/安全遗漏/低置信度 → Agent收到反馈后修订 (最多3轮)"),
    ("步骤 6-8: 仲裁决断阶段", ACCENT2,
     "SafetyPanel 独立规则审计(与辩论并行) → 安全底线不可破\n"
     "ModeratorAgent 合成所有轮次辩论 → 结构化总结\n"
     "ChiefReviewer.arbitrate() → 规则优先 + LLM意见融合 → rule_strict时硬证据不可覆盖"),
    ("步骤 9-10: 收尾阶段", ACCENT5,
     "CoordinatorAgent → 如需澄清则生成追问 (缺失字段/模糊信息)\n"
     "CaseStore → 完整CaseLog JSON持久化到 datasets/cases/\n"
     "MultiConsultResponse → 所有中间结果一次性返回前端"),
]

for i, (title, color, desc) in enumerate(detail_cards):
    y = Inches(2.5) + i * Inches(1.15)
    box = make_rounded_rect(slide, Inches(0.6), y, Inches(12), Inches(1.05), fill_color=CARD_BG, border_color=color)
    tf = box.text_frame; tf.word_wrap = True
    make_title_style(tf, title, size=16, color=color)
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line.strip()
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY

add_number(slide, 8, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
# SLIDE 9: 数据流全景图
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
            "数据流全景 — 输入到输出完整链路", size=32, color=WHITE, bold=True)

# Left column: Text pipeline
left_x = Inches(0.6)
box_l = make_rounded_rect(slide, left_x, Inches(1.2), Inches(5.8), Inches(5.8), fill_color=CARD_BG, border_color=ACCENT)
tf_l = box_l.text_frame; tf_l.word_wrap = True
make_title_style(tf_l, "文本药物审查数据流", size=18, color=ACCENT)
text_flow = [
    "",
    "INPUT: 临床文本 或 PatientContext + 候选药物",
    "   ↓",
    "[ExtractAgent] LLM解析 → 结构化患者信息",
    "   ↓  + Med7 NER辅助",
    "[ReviewEngine] 确定性规则 → 6维审查输出",
    "   ↓  + DDI-BERT分类",
    "[Agent Panel] 5专家并行LLM审查",
    "   ├─ ClinicalPharmacistAgent",
    "   ├─ InternalMedicineAgent",
    "   ├─ AllergySpecialistAgent",
    "   ├─ PharmacyInventoryAgent",
    "   └─ SpecialistAgent (条件触发)",
    "   ↓",
    "[Debate Engine] Critic ↔ Agents (≤3轮)",
    "   ↓  + SafetyPanel并行",
    "[Moderator → ChiefReviewer → Coordinator]",
    "   ↓",
    "OUTPUT: MultiConsultResponse (含全部中间结果)",
]
for item in text_flow:
    p = tf_l.add_paragraph()
    p.text = item
    p.font.size = Pt(11)
    if item.startswith("INPUT") or item.startswith("OUTPUT"):
        p.font.color.rgb = ACCENT5
        p.font.bold = True
    elif item.startswith("[") and "]" in item:
        p.font.color.rgb = ACCENT
    else:
        p.font.color.rgb = LIGHT_GRAY

# Right column: Imaging pipeline
right_x = Inches(6.7)
box_r = make_rounded_rect(slide, right_x, Inches(1.2), Inches(5.8), Inches(5.8), fill_color=CARD_BG, border_color=ACCENT3)
tf_r = box_r.text_frame; tf_r.word_wrap = True
make_title_style(tf_r, "影像审查数据流 (独立触发)", size=18, color=ACCENT3)
img_flow = [
    "",
    "INPUT: 影像选择 + 患者摘要",
    "   ↓",
    "[ImagingCatalog] 数据源扫描 → 列出可用研究",
    "   ↓  ImagingScope权限过滤",
    "[SegmentService] 串行加载分割模型",
    "   ├─ CXR-Lesion / BraTS-Tumor",
    "   ├─ TotalSegmentator / VISTA3D",
    "   └─ SAM-Med3D / SAM2D",
    "   ↓  分割标注叠加",
    "[OpenAIVisionClient] Qwen3-VL 多模态分析",
    "   ↓  base64图片 + 临床摘要",
    "VLM 结构化JSON输出",
    "   ↓",
    "[DeepSeekSynthesis] 融合 VLM结果 + Agent审查",
    "   ↓",
    "OUTPUT: ClinicalReport (7节结构化报告)",
]
for item in img_flow:
    p = tf_r.add_paragraph()
    p.text = item
    p.font.size = Pt(11)
    if item.startswith("INPUT") or item.startswith("OUTPUT"):
        p.font.color.rgb = ACCENT5
        p.font.bold = True
    elif item.startswith("[") and "]" in item:
        p.font.color.rgb = ACCENT3
    else:
        p.font.color.rgb = LIGHT_GRAY

add_number(slide, 9, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════
# SLIDE 10: 技术亮点 & 总结
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
            "技术亮点 & 设计原则", size=32, color=WHITE, bold=True)

highlights = [
    ("1. 安全优先的混合架构", ACCENT4,
     "确定性规则引擎 + LLM智能审查，规则作为不可覆盖的硬约束，rule_strict=true 时LLM无法推翻规则判定"),
    ("2. 模型分层解耦", ACCENT,
     "主LLM / 视觉LLM / 合成LLM / 传统ML 四层独立，每层独立配置、独立错误处理，可独立替换升级"),
    ("3. 多Agent对抗式协作", ACCENT3,
     "MDAgents模式的多轮辩论 → Critic挑战 → 修订收敛 → 首席仲裁，通过分歧驱动质量提升"),
    ("4. 影像处理串行优化", ACCENT2,
     "6个分割模型串行加载/释放GPU显存，VLM支持最多12张图片+分割标注叠加，减少显存压力"),
    ("5. 部门级权限隔离", ACCENT5,
     "ImagingScope 按科室控制影像源访问，Agent注册表支持科室特定专家动态激活，case_template按部门筛选"),
    ("6. 完整的可追溯性", WHITE,
     "CaseLog全流程记录(提取→规则→审查→辩论→仲裁→澄清→最终)，JSON持久化，便于审计和回溯"),
]

for i, (title, color, desc) in enumerate(highlights):
    y = Inches(1.2) + i * Inches(0.95)
    box = make_rounded_rect(slide, Inches(0.6), y, Inches(12), Inches(0.85), fill_color=CARD_BG, border_color=color)
    tf = box.text_frame; tf.word_wrap = True
    make_title_style(tf, title, size=16, color=color)
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY

# Bottom summary
box_b = make_rounded_rect(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.6), fill_color=RGBColor(0x10, 0x30, 0x40), border_color=ACCENT)
tf_b = box_b.text_frame; tf_b.word_wrap = True
p = tf_b.paragraphs[0]
p.text = "FastAPI  →  Orchestrator  →  Rule Engine + LLM Agents + Debate  →  VLM + Segmentation  →  DeepSeek Synthesis  →  Structured Report"
p.font.size = Pt(13)
p.font.color.rgb = ACCENT
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

add_number(slide, 10, TOTAL_SLIDES)

# ── Save ──
output_path = os.path.join(os.path.dirname(__file__), "..", "MedSafe_Backend_Architecture.pptx")
prs.save(output_path)
print(f"PPT saved to: {os.path.abspath(output_path)}")